from app import app
from flask import render_template, request, redirect, url_for
from pptx import Presentation
from pptx.util import Inches, Pt
import openai
import os
import re
from dotenv import load_dotenv
load_dotenv()
openai.api_key = os.environ.get('API_KEY')
import requests
import json

@app.route('/')
def index():

    return render_template('index.html')


@app.route('/create_ppt', methods=['POST'])
def create_ppt():
    def first_response_type(text):
        pattern = r"\d+\)\s"
        text = re.sub(pattern, "", text)
        text = text.replace("\n\n", "\n")
        pattern = r'Slide [Hh]eader (\d+): (.+)\n\d+\.\s(.+)\n\d+\.\s(.+)\n\d+\.\s(.+)\n'
        headers_and_content = re.findall(pattern, text)
        print(headers_and_content)
        if not headers_and_content:
            print("I am inside")
            pattern = r'Slide (\d+) [Hh]eader: (.+)\n\d+\.\s(.+)\n\d+\.\s(.+)\n\d+\.\s(.+)\n'
            headers_and_content = re.findall(pattern, text)
            print(headers_and_content)

        headers_and_content_dict = {}
        for header_and_content in headers_and_content:
            headers_and_content_dict[header_and_content[1]
                                     ] = header_and_content[2:]
        return headers_and_content_dict
    messages = [
        {"role": "system", "content": "You’re a kind helpful assistant"}
    ]

    def generate_header(prompt):

        messages.append({"role": "user", "content": prompt})
        completion = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=messages
        )
        return completion.choices[0].message.content

    topic = request.form['topic']
    slide_1_prompt = f"Please provide five slide header on {topic} in the following format: \n 1) Slide header 1: 1. Key Point 1, \n 2. Key Point 2, ,\n 3. Key Point 3 \n 2) Slide header 2: 1. Key Point 1, \n 2. Key Point 2, ,\n 3. Key Point 3 \n 3) Slide header 3: 1. Key Point 1, \n 2. Key Point 2, ,\n 3. Key Point 3 \n 4) Slide header 4: 1. Key Point 1, \n 2. Key Point 2, ,\n 3. Key Point 3 \n 5) Slide header 5: 1. Key Point 1, \n 2. Key Point 2, ,\n 3. Key Point 3"

    result_from_gpt = generate_header(slide_1_prompt)
    print(result_from_gpt)
    formatted_dict = first_response_type(result_from_gpt)
    print(formatted_dict)

    X = Presentation()

    Layout = X.slide_layouts[0]
    first_slide = X.slides.add_slide(Layout)

    first_slide.shapes.title.text = topic
    first_slide.placeholders[1].text = "Created by Big Data Lab"

    X.save("First_presentation.pptx")
    for i in formatted_dict:
        Layout = X.slide_layouts[5]
        slide = X.slides.add_slide(Layout)
        slide.shapes.title.text = i
        left = Inches(1)
        top = Inches(2)
        width = Inches(6)
        height = Inches(2)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        for bullet in formatted_dict[i]:
            p = textbox.text_frame.add_paragraph()
            p.text = bullet
            p.font.size = Pt(16)
            p.font.name = 'Calibri'
            p.level = 0
            p.font.bold = False
            # p.font.color.rgb = (0, 0, 0)
        # wrap the text within the textbox
        textbox.text_frame.word_wrap = True
        X.save("First_presentation.pptx")
    return redirect(url_for('index'))

