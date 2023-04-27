from pptx import Presentation
# from pptx.enum.text import PP_PARAGRAPH_STYLE
from pptx.util import Inches, Pt

# create a new presentation with one slide
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])

# define the bullet points as a list
bullet_points = ['point 1', 'point 2', 'point 3']

# create a text box shape on the slide
left = Inches(1)
top = Inches(2)
width = Inches(6)
height = Inches(2)
textbox = slide.shapes.add_textbox(left, top, width, height)

# add the bullet points to the text box as paragraphs
for point in bullet_points:
    p = textbox.text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.name = 'Calibri'
    p.level = 0
    p.font.bold = False
    # p.font.color.rgb = (0, 0, 0)
    bullet = p._p.add_bullet()
    bullet.char = u'\u2022'
    bullet.font.name = 'Wingdings'
    bullet.font.size = Pt(18)

# save the presentation
prs.save('my_presentation.pptx')